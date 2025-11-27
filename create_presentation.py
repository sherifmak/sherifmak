from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import io

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Swiss minimalist color palette
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
ACCENT = RGBColor(220, 53, 69)  # Red accent for Bainkom brand

def create_blank_slide(prs):
    """Create a blank slide"""
    blank_slide_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_slide_layout)

def add_title_slide(prs):
    """Create title slide"""
    slide = create_blank_slide(prs)

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "LATEST AI ADVANCEMENTS"
    title_para = title_frame.paragraphs[0]
    title_para.font.name = "Helvetica"
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = BLACK
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "For Interior Design Innovation"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.name = "Helvetica"
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = BLACK
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Bainkom branding
    brand_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.5))
    brand_frame = brand_box.text_frame
    brand_frame.text = "BAINKOM"
    brand_para = brand_frame.paragraphs[0]
    brand_para.font.name = "Helvetica"
    brand_para.font.size = Pt(36)
    brand_para.font.bold = True
    brand_para.font.color.rgb = ACCENT
    brand_para.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, description, use_case, number):
    """Create content slide with Swiss minimalist design"""
    slide = create_blank_slide(prs)

    # Slide number (top left corner)
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(1), Inches(0.3))
    num_frame = num_box.text_frame
    num_frame.text = f"{number:02d}"
    num_para = num_frame.paragraphs[0]
    num_para.font.name = "Helvetica"
    num_para.font.size = Pt(18)
    num_para.font.color.rgb = ACCENT

    # Title (top)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.name = "Helvetica"
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = BLACK

    # Description (left column)
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(2.5))
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True
    desc_frame.text = description
    desc_para = desc_frame.paragraphs[0]
    desc_para.font.name = "Helvetica"
    desc_para.font.size = Pt(14)
    desc_para.font.color.rgb = BLACK
    desc_para.line_spacing = 1.3

    # Use case label
    label_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(4.5), Inches(0.3))
    label_frame = label_box.text_frame
    label_frame.text = "USE CASE"
    label_para = label_frame.paragraphs[0]
    label_para.font.name = "Helvetica"
    label_para.font.size = Pt(11)
    label_para.font.bold = True
    label_para.font.color.rgb = ACCENT

    # Use case (left column bottom)
    use_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.9), Inches(4.5), Inches(2.2))
    use_frame = use_box.text_frame
    use_frame.word_wrap = True
    use_frame.text = use_case
    use_para = use_frame.paragraphs[0]
    use_para.font.name = "Helvetica"
    use_para.font.size = Pt(13)
    use_para.font.color.rgb = BLACK
    use_para.line_spacing = 1.3

    # Image placeholder (right side) - simple minimalist rectangle
    img_left = Inches(5.3)
    img_top = Inches(1.8)
    img_width = Inches(4.2)
    img_height = Inches(5.2)

    # Create a simple placeholder image with PIL
    img = Image.new('RGB', (800, 990), color=(245, 245, 245))
    draw = ImageDraw.Draw(img)

    # Draw a simple cross/plus symbol in the center (Swiss style)
    center_x, center_y = 400, 495
    line_length = 60
    line_width = 2
    draw.line([(center_x - line_length, center_y), (center_x + line_length, center_y)],
              fill=(200, 200, 200), width=line_width)
    draw.line([(center_x, center_y - line_length), (center_x, center_y + line_length)],
              fill=(200, 200, 200), width=line_width)

    # Save to bytes
    img_byte_arr = io.BytesIO()
    img.save(img_byte_arr, format='PNG')
    img_byte_arr.seek(0)

    # Add image to slide
    slide.shapes.add_picture(img_byte_arr, img_left, img_top, width=img_width, height=img_height)

# AI Advancements data
advancements = [
    {
        "title": "Google Gemini 3 Pro & Antigravity",
        "description": "Google's most intelligent model featuring state-of-the-art reasoning with a 1501 Elo score on LMArena. Integrated into Google Search on launch day. Antigravity is their new agentic development platform enabling developers to work at a task-oriented level with autonomous agents managing editor, terminal, and browser simultaneously.",
        "use_case": "Enterprise coding workflows where developers need AI agents to autonomously plan, code, test, and validate applications across multiple tools. For Bainkom: Automate complex backend development, automatically update interior design matching algorithms, and coordinate multi-system integrations across your marketplace platform."
    },
    {
        "title": "Claude Opus 4.5",
        "description": "Anthropic's newest flagship model described as 'the best model in the world for coding, agents, and computer use,' achieving 80.9% on SWE-bench Verified benchmark. Pricing reduced to $5/$25 per million tokens, making it significantly more accessible.",
        "use_case": "Long-horizon autonomous coding tasks, complex multi-system debugging, and agentic workflows requiring sustained reasoning across 30-minute sessions. For Bainkom: Build intelligent chatbots for customer-designer matching, automate code reviews, and create AI agents that handle customer onboarding workflows end-to-end."
    },
    {
        "title": "OpenAI AgentKit & Agent Builder",
        "description": "Complete toolkit featuring Agent Builder (visual canvas for multi-agent workflows), Connector Registry, ChatKit, and expanded evaluation capabilities. Agent Builder described as 'like Canva for building agents' - a fast, visual way to design agent logic.",
        "use_case": "Rapid agent prototyping with visual drag-and-drop interface, reducing complex agent orchestration from months to hours. For Bainkom: Quickly build multi-step workflows for designer vetting, customer preference collection, and automated project matching without extensive coding."
    },
    {
        "title": "LangChain Deep Agents Framework",
        "description": "Framework for building autonomous agents with planning tools, sub-agents, file system access, and detailed prompts inspired by Claude Code and Deep Research applications. Recent integration with Anthropic's skill-based framework enhances AI performance with dynamic skill folders.",
        "use_case": "Building research agents that can conduct multi-step investigations, maintain context across long tasks, and spawn specialized sub-agents for focused work. For Bainkom: Create agents that research interior design trends, analyze designer portfolios, and generate market insights from competitor platforms."
    },
    {
        "title": "OpenAI Agents SDK",
        "description": "Production-ready SDK for orchestrating multi-agent workflows with easily configurable agents, intelligent handoffs, guardrails, and built-in tracing. Replaced the experimental Swarm framework.",
        "use_case": "Production-scale multi-agent systems requiring robust orchestration, state management, and observability across complex handoff patterns. For Bainkom: Coordinate multiple AI agents handling customer inquiries, designer scheduling, project estimation, and payment processing with seamless handoffs."
    },
    {
        "title": "Nano Banana Pro (Gemini 3 Pro Image)",
        "description": "Google's state-of-the-art image generation and editing model built on Gemini 3 Pro, using advanced reasoning and real-world knowledge for superior visualization. Supports 2K and 4K image generation with exceptional text rendering in multiple languages, fonts, and styles.",
        "use_case": "Creating marketing materials with accurate multilingual text overlays, generating educational infographics with real-time data, and producing print-quality designs at 4K resolution. For Bainkom: Generate high-quality mood boards, create personalized design concept visualizations for clients, and produce marketing materials showcasing designer portfolios."
    },
    {
        "title": "Figma Weave (formerly Weavy.ai)",
        "description": "Node-based AI platform bringing generative AI and professional editing tools together, acquired by Figma to build image, video, animation, motion design, and VFX capabilities. Integrates multiple AI models including Sora, Veo for video, and Flux, Ideogram, Nano-Banana for images with granular editing controls.",
        "use_case": "Professional creative workflows combining AI generation with precise editing—architects staging spaces, VFX artists creating game assets, marketers producing social media content with full control over lighting, masking, and color grading. For Bainkom: Create realistic room staging, generate before/after visualizations, and produce video tours of design concepts."
    },
    {
        "title": "FLUX.1 Krea [dev] with ComfyUI",
        "description": "Open-weights text-to-image model from Black Forest Labs and Krea AI delivering realistic, diverse images that avoid the 'AI look' with natural details and exceptional quality. Recently integrated into ComfyUI with up to 40% performance improvements for NVIDIA RTX GPUs.",
        "use_case": "Node-based image generation workflows requiring photorealistic outputs without oversaturation—ideal for editorial photography, product visualization, and cinematic storyboarding with full workflow customization. For Bainkom: Generate photorealistic interior renders, create product visualizations for furniture placement, and produce portfolio-quality design mockups."
    },
    {
        "title": "Runway Gen-4",
        "description": "Runway's highest-fidelity AI video generator enabling consistent characters, locations, and objects across scenes with coherent world environments and realistic motion. Introduces character and scene consistency across multiple shots, solving AI video's biggest challenge through persistent memory of visual elements.",
        "use_case": "Cinematic storytelling with recurring characters across different lighting conditions and camera angles, VFX production for film/TV, and product visualization from multiple perspectives while maintaining brand consistency. For Bainkom: Create walkthrough videos of design concepts, showcase designer portfolios with consistent branding, and generate promotional videos for marketplace listings."
    },
    {
        "title": "Krea AI Node-Based Workflows",
        "description": "Platform providing access to 50+ models across images, videos, audio, and 3D on an infinite canvas with advanced controls for lighting, lens effects, object swapping, and color grading.",
        "use_case": "Complex creative pipelines chaining multiple AI models together—from initial image generation through video transformation to final color grading—with reusable community templates for rapid iteration. For Bainkom: Build end-to-end design visualization pipelines, create customizable templates for different room types, and automate style transfer across design concepts."
    },
    {
        "title": "n8n AI Workflow Enhancements",
        "description": "Platform evolution featuring agentic systems on single screen, visual workflows for agent interactions, 600+ AI automation templates, and Model Context Protocol (MCP) server integration.",
        "use_case": "Enterprise AI automation connecting LLMs with business systems—automating customer support workflows, content generation pipelines, and human-in-the-loop approval processes with compliance requirements. For Bainkom: Automate designer onboarding, sync customer preferences across systems, generate personalized email campaigns, and create approval workflows for project milestones."
    },
    {
        "title": "Meta SAM 3 & SAM 3D",
        "description": "Expansion of Segment Anything suite with SAM 3 for advanced image/video segmentation from text prompts, and SAM 3D for 3D reconstruction from 2D inputs.",
        "use_case": "E-commerce AR try-ons, automated product cutout creation for marketplaces, and 3D asset generation for gaming from 2D product photography. For Bainkom: Enable virtual furniture placement in customer photos, automatically extract furniture items from design inspiration images, and create 3D room models from 2D floor plans."
    }
]

# Add title slide
add_title_slide(prs)

# Add content slides
for i, adv in enumerate(advancements, 1):
    add_content_slide(prs, adv["title"], adv["description"], adv["use_case"], i)

# Save presentation
output_path = "/home/user/sherifmak/Bainkom_AI_Presentation.pptx"
prs.save(output_path)
print(f"Presentation created successfully: {output_path}")
